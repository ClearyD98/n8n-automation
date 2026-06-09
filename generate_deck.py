#!/usr/bin/env python3
"""autoprod deck generator — creates professional slide decks for client meetings"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, json, sys
from datetime import datetime

PURPLE = RGBColor(0x7C, 0x3A, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x0A, 0x0A, 0x0A)
MUTED = RGBColor(0x8A, 0x8A, 0x8E)
BORDER = RGBColor(0x1A, 0x1A, 0x1E)
W = Inches(13.333)
H = Inches(7.5)

class DeckBuilder:
    def __init__(self, client_name, service_type, output_dir=None):
        self.prs = Presentation()
        self.prs.slide_width = W
        self.prs.slide_height = H
        self.client = client_name
        self.service = service_type
        self.slides = []
        self.out = output_dir or os.path.expanduser("~/Projects/autoprod-relaunch/decks")
        os.makedirs(self.out, exist_ok=True)

    def _bg(self, slide, color=DARK):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_textbox(self, slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Inter'):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        return tf

    def title_slide(self, subtitle=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide)
        self._add_textbox(slide, 1.5, 2.5, 10, 1.5, f"autoprod × {self.client}", font_size=44, bold=True)
        title = {"audit": "AI Readiness Audit", "strategy": "AI Strategy Session", "retainer": "Monthly Advisory Review"}
        self._add_textbox(slide, 1.5, 4.0, 10, 1, title.get(self.service, self.service), font_size=28, color=MUTED)
        if subtitle:
            self._add_textbox(slide, 1.5, 4.8, 10, 0.8, subtitle, font_size=16, color=MUTED)
        self._add_textbox(slide, 1.5, 6.5, 10, 0.5, f"Dylan Cleary · dylan@autoprod.io · Dublin Docklands", font_size=12, color=MUTED)
        self.slides.append("title")

    def section_header(self, number, title, body=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide)
        self._add_textbox(slide, 1.5, 0.5, 1, 0.5, f"0{number}", font_size=16, color=PURPLE, bold=True)
        self._add_textbox(slide, 1.5, 1.5, 10, 1.5, title, font_size=36, bold=True)
        if body:
            self._add_textbox(slide, 1.5, 3.2, 8, 2, body, font_size=18, color=MUTED)
        self.slides.append(f"section_{number}")

    def bullet_slide(self, title, bullets):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide)
        self._add_textbox(slide, 1.5, 0.8, 10, 1, title, font_size=28, bold=True)
        y = 2.0
        for b in bullets:
            self._add_textbox(slide, 2.0, y, 9, 0.6, f"▸ {b}", font_size=16, color=WHITE)
            y += 0.55
        self.slides.append("bullets")

    def two_column(self, title, left_title, left_bullets, right_title, right_bullets):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide)
        self._add_textbox(slide, 1.5, 0.5, 10, 0.8, title, font_size=28, bold=True)
        # Left
        self._add_textbox(slide, 1.5, 1.8, 4.5, 0.5, left_title, font_size=20, bold=True, color=PURPLE)
        y = 2.5
        for b in left_bullets:
            self._add_textbox(slide, 1.8, y, 4.2, 0.5, f"▸ {b}", font_size=14, color=WHITE)
            y += 0.45
        # Right
        self._add_textbox(slide, 7.0, 1.8, 4.5, 0.5, right_title, font_size=20, bold=True, color=PURPLE)
        y = 2.5
        for b in right_bullets:
            self._add_textbox(slide, 7.3, y, 4.2, 0.5, f"▸ {b}", font_size=14, color=WHITE)
            y += 0.45
        self.slides.append("two_column")

    def next_steps(self, steps):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide)
        self._add_textbox(slide, 1.5, 1.0, 10, 1, "Next steps", font_size=36, bold=True)
        y = 2.5
        for i, step in enumerate(steps, 1):
            self._add_textbox(slide, 2.0, y, 9, 0.8, f"{i}.  {step}", font_size=20, color=WHITE)
            y += 1.2
        self.slides.append("next_steps")

    def contact_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide)
        self._add_textbox(slide, 1.5, 2.0, 10, 1.5, "Let's build this.", font_size=48, bold=True)
        self._add_textbox(slide, 1.5, 4.0, 10, 1, "Dylan Cleary", font_size=24, color=MUTED)
        self._add_textbox(slide, 1.5, 4.6, 10, 0.6, "dylan@autoprod.io  ·  Dublin Docklands", font_size=16, color=MUTED)
        self.slides.append("contact")

    def save(self, filename=None):
        if not filename:
            ts = datetime.now().strftime("%Y%m%d")
            filename = f"{self.client.lower().replace(' ','_')}_{self.service}_{ts}.pptx"
        path = os.path.join(self.out, filename)
        self.prs.save(path)
        print(f"✓ Deck saved: {path}")
        print(f"  Slides: {len(self.slides)} ({', '.join(self.slides)})")
        return path


def build_audit_deck(client_name, industry=""):
    """AI Readiness Audit deck"""
    d = DeckBuilder(client_name, "audit")
    d.title_slide(f"Prepared for {client_name}{' — ' + industry if industry else ''}  ·  {datetime.now().strftime('%B %Y')}")
    d.section_header(1, "Where you are today",
        "Every company is somewhere on the AI maturity curve. We'll map your current position — what's working, what's experimental, what's not even on the radar yet.")
    d.section_header(2, "What the EU AI Act means for you",
        "Not every AI system is regulated the same way. We classify your current and planned AI use against the Act's risk categories, so you know exactly what compliance looks like for your business.")
    d.section_header(3, "What's worth automating",
        "Not everything that can be automated should be. We identify the highest-ROI opportunities — the ones that save money, speed up operations, or give you a competitive edge — and rank them by impact.")
    d.bullet_slide("What we assessed", [
        "Current AI tools and systems in use",
        "Data readiness — what you have, what's missing",
        "Team capability — who knows what about AI",
        "Compliance posture against EU AI Act requirements",
        "Competitor AI activity in your sector",
        "Automation opportunities ranked by ROI",
    ])
    d.two_column("Key findings", "Strengths", [
        "Existing infrastructure that supports AI",
        "Data assets that can be leveraged",
        "Team members with AI aptitude",
    ], "Opportunities", [
        "Immediate automation wins (low effort, high impact)",
        "Compliance gaps to close before enforcement",
        "Quick wins to build momentum",
    ])
    d.section_header(4, "Your 90-day roadmap",
        "Week-by-week plan. What to do, in what order, with what resources. Prioritized so you get wins early while building toward larger transformation.")
    d.next_steps([
        "Review this deck with your leadership team",
        "Decide which automation opportunities to pursue first",
        "Book a follow-up strategy session to design the implementation plan",
        "Share the compliance checklist with your legal/compliance team",
    ])
    d.contact_slide()
    return d.save()


def build_strategy_deck(client_name, industry=""):
    """AI Strategy Session deck"""
    d = DeckBuilder(client_name, "strategy")
    d.title_slide(f"Prepared for {client_name}{' — ' + industry if industry else ''}  ·  {datetime.now().strftime('%B %Y')}")
    d.section_header(1, "Your AI ambition",
        "Where you want to be in 12 months. We define the vision — not 'AI for everything,' but specific, measurable outcomes that move your business forward.")
    d.section_header(2, "The automation roadmap",
        "A prioritized list of what to build, buy, or ignore. Each opportunity scored by: impact, cost, complexity, and compliance risk. You leave knowing exactly what to tackle first.")
    d.bullet_slide("Automation priorities", [
        "Tier 1: Immediate wins — 0-30 days, low effort, high impact",
        "Tier 2: Build phase — 30-90 days, requires some investment",
        "Tier 3: Transform — 90+ days, changes how you operate",
        "Tier 4: Avoid — looks good on a slide, terrible ROI in reality",
    ])
    d.section_header(3, "Compliance architecture",
        "How your AI systems fit within the EU AI Act. Risk classification for each system. Documentation requirements. Human oversight protocols. What needs to happen before enforcement and what can wait.")
    d.bullet_slide("Compliance requirements", [
        "Risk classification for each planned AI system",
        "Documentation and record-keeping obligations",
        "Human oversight requirements per risk tier",
        "Transparency obligations (user notification, labeling)",
        "Conformity assessment timeline",
        "Ongoing monitoring and reporting",
    ])
    d.two_column("Build vs Buy", "Build internally", [
        "Custom agentic workflows",
        "Internal LLM interfaces",
        "Proprietary data processing",
    ], "Buy / integrate", [
        "Off-the-shelf AI tools",
        "API-based AI services",
        "Platform-native AI features",
    ])
    d.section_header(4, "Implementation plan",
        "Team, timeline, budget. Who does what, when it ships, what success looks like. Enough detail to take to your board or investors.")
    d.next_steps([
        "Approve the automation roadmap with leadership",
        "Assign internal owner for AI implementation",
        "Begin compliance documentation for high-risk systems",
        "Schedule monthly retainer check-ins to track progress",
    ])
    d.contact_slide()
    return d.save()


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 generate_deck.py <client_name> <audit|strategy> [industry]")
        print("Example: python3 generate_deck.py 'FRS Recruitment' audit Recruitment")
        sys.exit(1)
    
    client = sys.argv[1]
    deck_type = sys.argv[2]
    industry = sys.argv[3] if len(sys.argv) > 3 else ""
    
    if deck_type == "audit":
        build_audit_deck(client, industry)
    elif deck_type == "strategy":
        build_strategy_deck(client, industry)
    else:
        print(f"Unknown deck type: {deck_type}")
        sys.exit(1)
